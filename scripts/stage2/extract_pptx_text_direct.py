import sys
from pathlib import Path

from pptx import Presentation

SKILL_DIR = Path(__file__).resolve().parent.parent.parent
sys.path.insert(0, str(SKILL_DIR))

from pptx_utils import configure_utf8_output


def extract_text(pptx_path: Path) -> str:
    prs = Presentation(str(pptx_path))
    lines: list[str] = []
    for slide_index, slide in enumerate(prs.slides, 1):
        lines.append(f"<!-- Slide number: {slide_index} -->")
        lines.append("")
        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                table = shape.table
                for row in table.rows:
                    row_text = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                    if any(row_text):
                        lines.append(" | ".join(row_text))
                lines.append("")
                continue

            if not getattr(shape, "has_text_frame", False):
                continue
            text = shape.text.strip()
            if text:
                lines.append(text)
                lines.append("")
    return "\n".join(lines).rstrip() + "\n"


def main() -> None:
    configure_utf8_output()
    if len(sys.argv) != 3:
        print("Usage: python extract_pptx_text_direct.py <input.pptx> <output.md>")
        sys.exit(1)

    pptx_path = Path(sys.argv[1])
    output_path = Path(sys.argv[2])
    if not pptx_path.exists():
        print(f"[error] PPTX not found: {pptx_path}")
        sys.exit(1)

    output_path.write_text(extract_text(pptx_path), encoding="utf-8")
    print(f"[done] Direct PPTX text: {output_path}")


if __name__ == "__main__":
    main()
